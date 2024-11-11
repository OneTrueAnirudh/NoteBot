#9/11 ver 2
from pptx import Presentation
from g4f.client import Client
from sentence_transformers import SentenceTransformer
import time
from collections import Counter

def read_pptx(file_path):
    """Reads a .pptx file and extracts text from all slides."""
    prs = Presentation(file_path)
    text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    return "\n".join(text)

def llm_keywords(slide_text):
    """Generates up to 15 unique keywords from slide text using an LLM and returns them as a list."""
    
    # Define the chunk size to prevent overloading the model
    chunk_size = 1000  # Adjust based on model's token limit
    
    # Split the slide_text into smaller chunks
    slide_chunks = [slide_text[i:i + chunk_size] for i in range(0, len(slide_text), chunk_size)]
    
    keywords = []
    
    # Initialize the client and handle errors
    client = Client()
    
    for chunk in slide_chunks:
        attempt = 0
        while attempt < 3:  # Retry up to 3 times per chunk
            try:
                prompt = "'{text}'\nread the above slide text, and give me 15 keywords/key phrases that describe the topics covered, for me to use as a query in my vector database that I want to synthesize notes with. respond only with the 15 keywords/phrases, each in 1 line, with no bullet points or numbering. also, avoid using brackets or short forms, just give me the pure keywords/phrases."
                
                response = client.chat.completions.create(
                    model="gpt-4o",  # Ensure the correct model name
                    messages=[{"role": "user", "content": prompt.format(text=chunk)}],
                )
                
                # Capture the response text and split into keywords
                result = response.choices[0].message.content
                keywords_list = result.split("\n")
                
                if "Model not found" in result or "too long input" in result or "error" in result:
                    raise ValueError("Error message received from the model.")
                
                # Add keywords from this chunk, stripping whitespace
                keywords.extend([keyword.strip() for keyword in keywords_list if keyword.strip()])
                break  # If successful, break out of the retry loop
            
            except Exception as e:
                print(f"Attempt {attempt + 1} failed for chunk: {e}")
                attempt += 1
                time.sleep(2)  # Wait before retrying
    
    if not keywords:
        print("Error: No keywords were generated.")
        return ["Model not found or too long input. Or any other error (xD)"]  # Fallback error message

    # Count occurrences of keywords and pick the top 15 unique keywords
    keyword_counts = Counter(keywords)
    top_keywords = [keyword for keyword, _ in keyword_counts.most_common(15)]
    
    return top_keywords



def get_keyword_vectors(keywords):
    """Converts a list of keywords into vectors using the all-mpnet-base-v2 model."""
    model = SentenceTransformer('sentence-transformers/all-mpnet-base-v2')
    keyword_vectors = model.encode(keywords)
    return keyword_vectors

# Example of usage:
# if __name__ == "__main__":
#     file_path = r"C:\Users\theon\Downloads\DM CAT-2\Google-SEO-Search-Engine-Optimization-Introduction-Powerpoint-Presentation-.pptx"
    
#     slide_text = read_pptx(file_path)
#     keywords = llm_keywords(slide_text)
#     print("Keywords Extracted:", keywords)
    
#     keyword_vectors = get_keyword_vectors(keywords)
#     print(f"Number of Keyword Vectors: {len(keyword_vectors)}")
