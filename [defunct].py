from transformers import pipeline, PegasusTokenizer, PegasusForConditionalGeneration
import spacy
from pptx import Presentation
from keybert import KeyBERT
import warnings

warnings.filterwarnings("ignore")  # suppresses all warnings

# Load the SpaCy model
nlp = spacy.load("en_core_web_sm")

# Load Pegasus model and tokenizer
tokenizer = PegasusTokenizer.from_pretrained("google/pegasus-large")
model = PegasusForConditionalGeneration.from_pretrained("google/pegasus-large")

# Create summarization pipeline using the loaded model and tokenizer
summarizer = pipeline("summarization", model=model, tokenizer=tokenizer)

def read_pptx(file_path):
    """Reads a .pptx file and extracts text from all slides."""
    prs = Presentation(file_path)
    text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    return "\n".join(text)

def advanced_preprocess(text):
    """Preprocess text using SpaCy."""
    doc = nlp(text)
    
    cleaned_tokens = []
    for token in doc:
        if not token.is_stop and not token.is_punct:
            cleaned_tokens.append(token.lemma_)  # Use lemma for base form

    # Join tokens back into a cleaned string
    cleaned_text = ' '.join(cleaned_tokens)
    return cleaned_text

def chunk_by_character_count(text, max_length=1024):
    """Chunks text based on a maximum character count."""
    chunks = []
    
    while len(text) > max_length:
        # Find the last space within the max_length limit
        split_index = text.rfind(' ', 0, max_length)
        if split_index == -1:  # If no space found, force split
            split_index = max_length
        
        # Append the chunk and slice the text
        chunks.append(text[:split_index].strip())
        text = text[split_index:].strip()
    
    if text:  # Append any remaining text
        chunks.append(text)
        
    return chunks

def summarize_text(text):
    """Summarizes text using Pegasus with chunking by character count."""
    chunks = chunk_by_character_count(text)
    summaries = []

    for chunk in chunks:
        if chunk.strip():  # Check for non-empty chunk
            summarized = summarizer(chunk, max_length=150, min_length=30, do_sample=False)
            summaries.append(summarized[0]['summary_text'])

    # Combine all summaries into a single summary
    final_summary = ' '.join(summaries)
    return final_summary

def extract_keywords_from_pptx(file_path):
    """Extract keywords from a .pptx file."""
    # Step 1: Read text from the PPTX file
    slide_text = read_pptx(file_path)

    # Step 2: Preprocess the extracted text
    cleaned_text = advanced_preprocess(slide_text)

    # Step 3: Summarize the cleaned text
    summarized_text = summarize_text(cleaned_text)
    print(summarized_text)

    # Step 4: Use KeyBERT to extract keywords
    keybert_model = KeyBERT()
    keywords = keybert_model.extract_keywords(summarized_text, keyphrase_ngram_range=(1, 1), top_n=15, diversity=0.3, use_mmr=True, stop_words='english')
    
    # Return just the keywords as a list
    return [keyword[0] for keyword in keywords]

# Sample usage
if __name__ == "__main__":
    file_path = r"C:\Users\theon\Downloads\DM CAT-2\Google-SEO-Search-Engine-Optimization-Introduction-Powerpoint-Presentation-.pptx"  # Replace with your PPTX file path
    keywords = extract_keywords_from_pptx(file_path)
    
    print("Extracted Keywords:")
    for keyword in keywords:
        print(keyword)
